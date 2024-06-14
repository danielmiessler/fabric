import argparse
import os
from urllib.parse import urljoin
from playwright.sync_api import sync_playwright
from bs4 import BeautifulSoup, Tag
from markdownify import MarkdownConverter, markdownify as md
import subprocess
import pdfplumber
import pandas as pd
from pptx import Presentation
from docx import Document


class CustomMarkdownConverter(MarkdownConverter):
    def __init__(self, *args, base_url=None, include_links=False, **kwargs):
        self.base_url = base_url
        self.links = []
        self.include_links = include_links
        super().__init__(*args, **kwargs)

    def convert_div(self, el, text, convert_as_inline):
        return text.strip() + "\n"

    def convert_span(self, el, text, convert_as_inline):
        return text

    def convert_a(self, el, text, convert_as_inline):
        href = el.get("href", "")
        if text.strip() == "" or href.startswith("#"):
            return ""
        absolute_url = urljoin(self.base_url, href) if self.base_url else href
        self.links.append((text, absolute_url))
        return ""

    def convert_img(self, el, text, convert_as_inline):
        src = el.get("src", "")
        alt = el.get("alt", "")
        if alt.strip() == "":
            return ""
        absolute_src = urljoin(self.base_url, src) if self.base_url else src
        self.links.append((alt, absolute_src))
        return ""

    def convert_section(self, el, text, convert_as_inline):
        return f"\n## {text.strip()}\n"

    def convert_article(self, el, text, convert_as_inline):
        return f"\n### {text.strip()}\n"

    def convert_header(self, el, text, convert_as_inline):
        if el.name in ["h1", "h2", "h3", "h4", "h5", "h6"]:
            level = int(el.name[1])
            return f'\n{"#" * level} {text.strip()}\n'
        return text

    def convert_footer(self, el, text, convert_as_inline):
        return f"\n---\n{text.strip()}\n---\n"

    def convert_nav(self, el, text, convert_as_inline):
        return f"\n### Navigation\n{text.strip()}\n"

    def convert_ul(self, el, text, convert_as_inline):
        return text.strip() + "\n"

    def convert_ol(self, el, text, convert_as_inline):
        return text.strip() + "\n"

    def convert_li(self, el, text, convert_as_inline):
        parent = el.parent
        if parent.name == "ul":
            text = text.strip()
            if not text:
                return ""
            else:
                return f"- {text.strip()}\n"
        elif parent.name == "ol":
            index = list(parent.children).index(el) + 1
            return f"{index}. {text.strip()}\n"
        return text.strip()

    def convert_strong(self, el, text, convert_as_inline):
        return f"**{text}**"

    def convert_em(self, el, text, convert_as_inline):
        return f"*{text}*"

    def convert_blockquote(self, el, text, convert_as_inline):
        return "> " + text.replace("\n", "\n> ").strip() + "\n"

    def convert_pre(self, el, text, convert_as_inline):
        return "\n```\n" + text.strip() + "\n```\n"

    def convert_code(self, el, text, convert_as_inline):
        if el.parent.name == "pre":
            return text  # 'pre'タグ内の'code'タグは'pre'タグで変換されるため無視
        return f"`{text}`"

    def convert_h1(self, el, text, convert_as_inline):
        return f"# {text.strip()}\n"

    def convert_h2(self, el, text, convert_as_inline):
        return f"## {text.strip()}\n"

    def convert_h3(self, el, text, convert_as_inline):
        return f"### {text.strip()}\n"

    def convert_h4(self, el, text, convert_as_inline):
        return f"#### {text.strip()}\n"

    def convert_h5(self, el, text, convert_as_inline):
        return f"##### {text.strip()}\n"

    def convert_h6(self, el, text, convert_as_inline):
        return f"###### {text.strip()}\n"

    def convert_del(self, el, text, convert_as_inline):
        return f"~~{text}~~"

    def convert_ins(self, el, text, convert_as_inline):
        return f"__{text}__"

    def convert_mark(self, el, text, convert_as_inline):
        return f"=={text}=="

    # Add more custom converters as needed


def fetch_page_content(url):
    with sync_playwright() as p:
        browser = p.firefox.launch(headless=True)
        context = browser.new_context(
            user_agent="Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
        )

        page = context.new_page()

        try:
            page.goto(url, timeout=20000)
            content = page.content()
        except Exception as e:
            content = curl_fetch(url)
        finally:
            browser.close()

        # 連続した改行を削除する
        content = "".join(filter(None, content.splitlines()))
        return content


def curl_fetch(url):
    try:
        result = subprocess.run(
            ["curl", "-s", url], capture_output=True, text=True, check=True
        )
        return result.stdout
    except subprocess.CalledProcessError as e:
        print(f"Error fetching {url} with curl: {e}")
        return ""


def parse_pdf(file_path):
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text()
    return text


def parse_docx(file_path):
    doc = Document(file_path)
    return "\n".join([p.text for p in doc.paragraphs])


def parse_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def parse_excel(file_path):
    df = pd.read_excel(file_path)
    return df.to_markdown(index=False)


def html_to_markdown(html_content, base_url=None, include_links=False):
    soup = BeautifulSoup(html_content, "html.parser")
    body = soup.find("body")
    if body:
        converter = CustomMarkdownConverter(
            base_url=base_url, include_links=include_links
        )
        markdown_content = converter.convert_soup(body).strip()
        if include_links and converter.links:
            markdown_content += "\n\n## Links\n"
            for text, url in converter.links:
                markdown_content += f"- [{text}]({url})\n"
        return markdown_content
    else:
        return "No body content found."


def main_function(url, options):
    if os.path.exists(url):
        file_path = url
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            text = parse_pdf(file_path)
        elif ext == ".docx":
            text = parse_docx(file_path)
        elif ext == ".pptx":
            text = parse_pptx(file_path)
        elif ext == ".xlsx" or ext == ".xls":
            text = parse_excel(file_path)
        else:
            with open(file_path, "r", encoding="utf-8") as file:
                text = file.read()
    else:
        html_content = fetch_page_content(url)
        text = html_to_markdown(
            html_content, base_url=url, include_links=options.include_links
        )
    print(text)


def main():
    parser = argparse.ArgumentParser(
        description="Fetch page content using Playwright and convert to Markdown."
    )
    parser.add_argument(
        "url", type=str, help="The URL of the page to fetch or the local file path."
    )
    parser.add_argument(
        "--include-links",
        action="store_true",
        help="Include links at the bottom of the markdown.",
    )
    args = parser.parse_args()

    if args.url is None:
        print("Error: No URL provided.")
        return

    main_function(args.url, args)


if __name__ == "__main__":
    main()
